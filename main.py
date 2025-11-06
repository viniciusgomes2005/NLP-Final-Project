import os
import argparse
from io import StringIO
from pathlib import Path
import numpy as np
import pandas as pd
import torch
import torch.nn as nn
from tqdm import tqdm
import matplotlib.pyplot as plt
from sklearn.metrics import accuracy_score, f1_score, classification_report
from sklearn.decomposition import PCA
import sentencepiece as spm

# =========================
# Configuração global
# =========================
OUT = Path("outputs")
(OUT / "figs").mkdir(parents=True, exist_ok=True)
PAD_TOKEN = "<PAD>"
DEVICE = torch.device("cuda" if torch.cuda.is_available() else "cpu")

# =========================
# Dados: IMDB (Hugging Face) com fallback offline
# =========================
def load_imdb(sample_train=2000, sample_test=1000, seed=42):
    try:
        from datasets import load_dataset
        ds = load_dataset("imdb")
        tr = ds["train"].shuffle(seed=seed).select(range(min(sample_train, len(ds["train"]))))
        te = ds["test"].shuffle(seed=seed).select(range(min(sample_test, len(ds["test"]))))
        X_train = list(tr["text"]); y_train = torch.tensor(tr["label"])
        X_test  = list(te["text"]);  y_test  = torch.tensor(te["label"])
        src = "imdb:hf"
    except Exception:
        # Fallback simples (sem internet)
        rng = np.random.default_rng(seed)
        pos = [
            "A lovely film with heart. Performances are touching and genuine.",
            "Brilliant direction, stellar pacing, and memorable characters.",
            "A moving story that kept me engaged; excellent cinematography.",
        ] * (max(1, sample_train // 6))
        neg = [
            "Terrible movie: dull plot and wooden acting.",
            "Painfully slow. I almost walked out halfway through.",
            "Messy writing and poor editing ruin any potential.",
        ] * (max(1, sample_train // 6))

        train_pairs = [(t,1) for t in pos[:sample_train//2]] + [(t,0) for t in neg[:sample_train//2]]
        rng.shuffle(train_pairs)
        X_train = [t for t,_ in train_pairs]
        y_train = torch.tensor([y for _,y in train_pairs])

        pos_t = pos[:sample_test//2]
        neg_t = neg[:sample_test//2]
        test_pairs = [(t,1) for t in pos_t] + [(t,0) for t in neg_t]
        rng.shuffle(test_pairs)
        X_test = [t for t,_ in test_pairs]
        y_test = torch.tensor([y for _,y in test_pairs])
        src = "fallback:mini"
    return X_train, y_train, X_test, y_test, src

# =========================
# SentencePiece
# =========================
def train_sentencepiece(corpus_texts, vocab_size, prefix, pad_token=PAD_TOKEN):
    input_data = "\n".join(corpus_texts)
    input_fp = StringIO(input_data)
    spm.SentencePieceTrainer.train(
        sentence_iterator=input_fp,
        model_prefix=prefix,
        vocab_size=vocab_size,
        user_defined_symbols=[pad_token],
    )
    sp = spm.SentencePieceProcessor()
    sp.load(f"{prefix}.model")
    return sp, sp.piece_to_id(pad_token)

# =========================
# Tokenização + Padding
# =========================
def pad_to_len(sequences, pad_idx, max_len):
    padded = []
    for s in sequences:
        if len(s) >= max_len:
            padded.append(s[:max_len])
        else:
            padded.append(s + [pad_idx] * (max_len - len(s)))
    return padded

def tokenize_batch(texts, sp, max_len, pad_idx):
    ids = [sp.encode_as_ids(t) for t in texts]
    ids = pad_to_len(ids, pad_idx, max_len)
    return torch.tensor(ids)

# =========================
# Modelos
# =========================
class SimpleClassifier(nn.Module):
    def __init__(self, vocab_size, embedding_dim, padding_idx):
        super().__init__()
        self.embedding = nn.Embedding(vocab_size, embedding_dim, padding_idx=padding_idx)
        self.clf = nn.Linear(embedding_dim, 1)

    def forward(self, x, return_embeddings=False):
        x = self.embedding(x)           # (B, L, D)
        doc_emb = torch.mean(x, dim=1)  # (B, D)
        logits = self.clf(doc_emb)      # (B, 1)
        if return_embeddings:
            return logits, doc_emb
        return logits

class GRUClassifier(nn.Module):
    def __init__(self, vocab_size, embedding_dim, padding_idx, hidden_dim=64):
        super().__init__()
        self.embedding = nn.Embedding(vocab_size, embedding_dim, padding_idx=padding_idx)
        self.gru = nn.GRU(input_size=embedding_dim, hidden_size=hidden_dim, batch_first=True)
        self.clf = nn.Linear(hidden_dim, 1)

    def forward(self, x, return_embeddings=False):
        emb = self.embedding(x)         # (B, L, D)
        out, h = self.gru(emb)          # h: (1, B, H)
        doc_emb = h.squeeze(0)          # (B, H)
        logits = self.clf(doc_emb)
        if return_embeddings:
            return logits, doc_emb
        return logits

# =========================
# Treino/Eval
# =========================
@torch.no_grad()
def evaluate(model, X, y, sp, max_len, pad_idx, pca_dims=2):
    model.eval()
    toks = tokenize_batch(X, sp, max_len, pad_idx).to(DEVICE)
    logits, emb = model(toks, return_embeddings=True)
    probs = torch.sigmoid(logits).cpu().numpy().flatten()
    preds = (probs > 0.5).astype(int)
    acc = accuracy_score(y.numpy(), preds)
    f1 = f1_score(y.numpy(), preds, average="macro", zero_division=0)
    rep = classification_report(y.numpy(), preds, zero_division=0, output_dict=True)
    emb_np = emb.cpu().numpy()
    if pca_dims and emb_np.shape[1] > pca_dims:
        emb_np = PCA(n_components=pca_dims, random_state=42).fit_transform(emb_np)
    return acc, f1, rep, emb_np, preds

def train_once(
    model, X_train, y_train, X_val, y_val,
    sp, pad_idx, max_len=100, epochs=25, lr=5e-3, batch_size=64, pca_dims=2
):
    model = model.to(DEVICE)
    opt = torch.optim.Adam(model.parameters(), lr=lr)
    criterion = nn.BCEWithLogitsLoss()

    train_losses, val_losses = [], []
    for ep in range(epochs):
        model.train()
        idx = np.random.permutation(len(X_train))
        for i in range(0, len(X_train), batch_size):
            sel = idx[i:i+batch_size]
            toks = tokenize_batch([X_train[j] for j in sel], sp, max_len, pad_idx).to(DEVICE)
            yb = y_train[sel].to(DEVICE).float().unsqueeze(1)

            opt.zero_grad()
            logits = model(toks)
            loss = criterion(logits, yb)
            loss.backward()
            opt.step()

        train_losses.append(loss.item())  # última loss do último batch

        model.eval()
        toks_val = tokenize_batch(X_val, sp, max_len, pad_idx).to(DEVICE)
        yv = y_val.to(DEVICE).float().unsqueeze(1)
        with torch.no_grad():
            logits_val = model(toks_val)
            vloss = criterion(logits_val, yv).item()
        val_losses.append(vloss)

    acc_tr, f1_tr, rep_tr, emb_tr, _ = evaluate(model, X_train, y_train, sp, max_len, pad_idx, pca_dims=pca_dims)
    acc_te, f1_te, rep_te, emb_te, preds = evaluate(model, X_val,   y_val,   sp, max_len, pad_idx, pca_dims=pca_dims)

    return {
        "model": model,
        "train_losses": train_losses,
        "val_losses": val_losses,
        "acc_train": acc_tr, "f1_train": f1_tr, "rep_train": rep_tr,
        "acc_test":  acc_te, "f1_test":  f1_te, "rep_test":  rep_te,
        "emb_train": emb_tr, "emb_test": emb_te,
        "preds": preds,
    }

# =========================
# Plot helpers
# =========================
def plot_losses(name, train_losses, val_losses):
    path = OUT / "figs" / f"{name}_loss.png"
    plt.figure(figsize=(5,3))
    plt.plot(train_losses, label="Train Loss")
    plt.plot(val_losses,   label="Eval  Loss")
    plt.xlabel("Epoch"); plt.ylabel("Loss"); plt.title(f"Loss — {name}")
    plt.legend(); plt.tight_layout(); plt.savefig(path); plt.close()
    return str(path)

def plot_embeddings(name, emb, labels, split):
    path = OUT / "figs" / f"{name}_emb_{split}.png"
    plt.figure(figsize=(4,3))
    sc = plt.scatter(emb[:,0], emb[:,1], c=np.array(labels), alpha=0.85)
    plt.xlabel("Dim 1"); plt.ylabel("Dim 2"); plt.title(f"{split.capitalize()} Embeddings — {name}")
    plt.tight_layout(); plt.savefig(path); plt.close()
    return str(path)

# =========================
# Experimentos
# =========================
def run_experiment(
    name, X_train, y_train, X_test, y_test,
    vocab_size=4000, embedding_dim=32, max_len=100, epochs=25, lr=5e-3,
    sample_train=None, sample_test=None, model_type="mean"
):
    # Subamostragem opcional
    X_tr = X_train[:sample_train] if sample_train else X_train
    y_tr = y_train[:sample_train] if sample_train else y_train
    X_te = X_test[:sample_test]   if sample_test   else X_test
    y_te = y_test[:sample_test]   if sample_test   else y_test

    # Treina SentencePiece específico do experimento
    sp, pad_idx = train_sentencepiece(X_tr, vocab_size=vocab_size, prefix=f"sp_{name}")

    # Modelo
    if model_type == "gru":
        model = GRUClassifier(vocab_size=len(sp), embedding_dim=embedding_dim, padding_idx=pad_idx)
    else:
        model = SimpleClassifier(vocab_size=len(sp), embedding_dim=embedding_dim, padding_idx=pad_idx)

    # Treino
    res = train_once(
        model, X_tr, y_tr, X_te, y_te, sp, pad_idx,
        max_len=max_len, epochs=epochs, lr=lr, batch_size=64, pca_dims=2
    )

    # Plots
    loss_plot = plot_losses(name, res["train_losses"], res["val_losses"])
    emb_train_plot = plot_embeddings(name, res["emb_train"], y_tr.numpy(), split="train")
    emb_test_plot  = plot_embeddings(name, res["emb_test"],  y_te.numpy(), split="test")

    summary = {
        "name": name,
        "vocab_size": vocab_size,
        "embedding_dim": embedding_dim,
        "max_len": max_len,
        "epochs": epochs,
        "model_type": model_type,
        "acc_train": res["acc_train"],
        "acc_test":  res["acc_test"],
        "f1_train":  res["f1_train"],
        "f1_test":   res["f1_test"],
        "loss_plot": loss_plot,
        "emb_train_plot": emb_train_plot,
        "emb_test_plot":  emb_test_plot,
    }
    return summary

# =========================
# Deck de slides (<= 5)
# =========================
def build_deck(df):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()

    def add_title_and_content(title, bullets):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.add_paragraph() if i>0 else tf.paragraphs[0]
            p.text = b
            p.level = 0

    def add_image_slide(title, img_path):
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        left = Inches(0.7); top = Inches(1.7); width = Inches(8.0)
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, left, top, width=width)

    # Slide 1 — Resumo
    best = df.loc[df["acc_test"].idxmax()]
    bullets = [
        "Pipeline: SentencePiece → Embedding → Média → Linear",
        "Treino: BCEWithLogits + Adam; curvas de loss + PCA(2D) para embeddings",
        "Parâmetros varridos: vocab_size, dataset_size, embedding_dim",
        f"Melhor experimento: {best['name']} (acc_test={best['acc_test']:.3f})",
    ]
    add_title_and_content("Resumo dos Experimentos", bullets)

    # Gráfico geral
    # salva figura de resumo
    plt.figure(figsize=(6,3))
    plt.plot(df["name"], df["acc_test"], marker="o")
    plt.xticks(rotation=30, ha="right")
    plt.ylabel("Accuracy (test)"); plt.title("Accuracy por experimento")
    plt.tight_layout()
    summary_fig = OUT / "figs" / "summary_accuracy.png"
    plt.savefig(summary_fig); plt.close()
    add_image_slide("Accuracy por experimento", str(summary_fig))

    # 2-4 — detalhes (vocab/data/emb)
    groups = [("vocab_", "Impacto do Vocabulário"),
              ("data_",  "Impacto do Tamanho do Dataset"),
              ("emb_",   "Impacto da Dimensão do Embedding")]
    for prefix, title in groups:
        cand = df[df["name"].str.startswith(prefix)]
        if not len(cand):
            continue
        b = cand.iloc[cand["acc_test"].values.argmax()]
        bullets = [
            f"Experimento: {b['name']}",
            f"vocab_size={b['vocab_size']} | emb_dim={b['embedding_dim']} | epochs={b['epochs']}",
            f"Accuracy={b['acc_test']:.3f} | F1={b['f1_test']:.3f}",
            "Abaixo: curvas de loss e separação no embedding (teste).",
        ]
        add_title_and_content(title, bullets)
        add_image_slide(f"{title} — Loss", b["loss_plot"])
        add_image_slide(f"{title} — Embeddings (test)", b["emb_test_plot"])

    # Limita a 5 slides
    while len(prs.slides) > 5:
        r = prs.slides._sldIdLst[-1]
        prs.slides._sldIdLst.remove(r)

    out_pptx = OUT / "slides_experimentos_embeddings.pptx"
    prs.save(out_pptx)
    return str(out_pptx)

# =========================
# Modo grid e quick
# =========================
def run_grid(model_type="mean", epochs=25, max_len=100, seed=42):
    X_train, y_train, X_test, y_test, src = load_imdb(sample_train=2000, sample_test=1000, seed=seed)
    print(f"[dados] origem={src} | train={len(X_train)} | test={len(X_test)}")

    configs = [
        # Vocabulário
        dict(name="vocab_2k", vocab_size=2000, embedding_dim=32, sample_train=2000, sample_test=1000),
        dict(name="vocab_8k", vocab_size=8000, embedding_dim=32, sample_train=2000, sample_test=1000),
        # Dataset size
        dict(name="data_1k",  vocab_size=4000, embedding_dim=32, sample_train=1000, sample_test=800),
        dict(name="data_3k",  vocab_size=4000, embedding_dim=32, sample_train=min(3000, len(X_train)), sample_test=1000),
        # Embedding dim
        dict(name="emb_16",   vocab_size=4000, embedding_dim=16, sample_train=2000, sample_test=1000),
        dict(name="emb_64",   vocab_size=4000, embedding_dim=64, sample_train=2000, sample_test=1000),
    ]

    summaries = []
    for cfg in configs:
        s = run_experiment(
            **cfg, X_train=X_train, y_train=y_train, X_test=X_test, y_test=y_test,
            max_len=max_len, epochs=epochs, model_type=model_type
        )
        summaries.append(s)

    df = pd.DataFrame(summaries)
    df.to_csv(OUT / "experiments_summary.csv", index=False)
    deck_path = build_deck(df)
    print(f"[ok] resumo: {OUT/'experiments_summary.csv'}")
    print(f"[ok] slides: {deck_path}")

def run_quick(model_type="mean", epochs=10, max_len=80, seed=42):
    X_train, y_train, X_test, y_test, src = load_imdb(sample_train=1000, sample_test=500, seed=seed)
    print(f"[dados] origem={src} | train={len(X_train)} | test={len(X_test)}")

    s = run_experiment(
        name="quick_demo", vocab_size=3000, embedding_dim=32,
        sample_train=1000, sample_test=500,
        X_train=X_train, y_train=y_train, X_test=X_test, y_test=y_test,
        max_len=max_len, epochs=epochs, model_type=model_type
    )
    df = pd.DataFrame([s])
    df.to_csv(OUT / "experiments_summary.csv", index=False)
    deck_path = build_deck(df)
    print(f"[ok] resumo: {OUT/'experiments_summary.csv'}")
    print(f"[ok] slides: {deck_path}")

# =========================
# CLI
# =========================
def main():
    ap = argparse.ArgumentParser(description="Experimentos NLP com Embeddings + Slides")
    ap.add_argument("--mode", choices=["grid", "quick"], default="grid", help="grid (todos) ou quick (rápido)")
    ap.add_argument("--model", choices=["mean", "gru"], default="mean", help="mean-pool (linear) ou GRU (extensão)")
    ap.add_argument("--epochs", type=int, default=25)
    ap.add_argument("--max-len", type=int, default=100)
    ap.add_argument("--seed", type=int, default=42)
    args = ap.parse_args()

    if args.mode == "grid":
        run_grid(model_type=args.model, epochs=args.epochs, max_len=args.max_len, seed=args.seed)
    else:
        run_quick(model_type=args.model, epochs=args.epochs, max_len=args.max_len, seed=args.seed)

if __name__ == "__main__":
    main()
